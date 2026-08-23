#!/usr/bin/env python3
"""
Builds the presentation for DS Project 2 (Hostel Student Management System).
Run:  python3 docs/make_ppt.py     ->  docs/DS_Project_2_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------- theme ----------
NAVY   = RGBColor(0x12, 0x35, 0x6F)
BLUE   = RGBColor(0x2F, 0x6F, 0xED)
GREEN  = RGBColor(0x16, 0xA3, 0x4A)
RED    = RGBColor(0xDC, 0x26, 0x26)
DARK   = RGBColor(0x1A, 0x1A, 0x1A)
GREY   = RGBColor(0x55, 0x5F, 0x6D)
LIGHT  = RGBColor(0xF3, 0xF6, 0xFB)
CODEBG = RGBColor(0x1E, 0x24, 0x30)
CODEFG = RGBColor(0xE6, 0xE9, 0xEF)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xD3, 0xDA, 0xE5)

BODY = "Calibri"
MONO = "Consolas"

W, H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width, prs.slide_height = W, H


def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])


def box(slide, x, y, w, h, fill=None, line=None, lw=1.25):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    return sh


def text(slide, x, y, w, h, runs, size=16, color=DARK, font=BODY,
         align=PP_ALIGN.LEFT, bold=False, space=6, line=None, anchor=MSO_ANCHOR.TOP):
    """runs = string, or list of paragraphs; a paragraph is a string or a list of
       (text, bold, color) tuples."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    paras = [runs] if isinstance(runs, str) else runs
    for i, p in enumerate(paras):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        para.space_after = Pt(space)
        if line:
            para.line_spacing = line
        pieces = [(p, bold, color)] if isinstance(p, str) else p
        for t, b, c in pieces:
            r = para.add_run()
            r.text = t
            r.font.size = Pt(size)
            r.font.bold = b
            r.font.color.rgb = c
            r.font.name = font
    return tb


def header(slide, title, kicker=None):
    box(slide, 0, 0, W, Inches(1.0), fill=NAVY)
    text(slide, Inches(0.55), Inches(0.24), Inches(11.0), Inches(0.55),
         title, size=27, color=WHITE, bold=True)
    if kicker:
        text(slide, Inches(11.3), Inches(0.34), Inches(1.6), Inches(0.4),
             kicker, size=13, color=RGBColor(0x9F, 0xB6, 0xE8), bold=True,
             align=PP_ALIGN.RIGHT)


def code(slide, x, y, w, lines, size=13.5, accent=BLUE):
    h = Inches(0.30 + 0.245 * len(lines))
    box(slide, x, y, w, h, fill=CODEBG)
    box(slide, x, y, Inches(0.055), h, fill=accent)
    tb = slide.shapes.add_textbox(x + Inches(0.22), y + Inches(0.14),
                                  w - Inches(0.34), h - Inches(0.2))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, ln in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.space_after = Pt(0)
        para.line_spacing = 1.12
        # anything after '//' is a comment -> green
        main, sep, com = ln.partition("//")
        r = para.add_run(); r.text = main
        r.font.size = Pt(size); r.font.name = MONO; r.font.color.rgb = CODEFG
        if sep:
            r2 = para.add_run(); r2.text = sep + com
            r2.font.size = Pt(size); r2.font.name = MONO
            r2.font.color.rgb = RGBColor(0x7E, 0xC9, 0x9B)
    return y + h


def bullets(slide, x, y, w, items, size=16.5, gap=10, marker=BLUE):
    """items = list of (bold_lead, rest) or plain strings."""
    cur = y
    for it in items:
        box(slide, x, cur + Inches(0.085), Inches(0.075), Inches(0.075), fill=marker)
        if isinstance(it, str):
            runs = [(it, False, DARK)]
        else:
            runs = [(it[0], True, NAVY), (it[1], False, DARK)]
        tb = text(slide, x + Inches(0.26), cur, w - Inches(0.26), Inches(0.4),
                  [runs], size=size, line=1.25)
        # estimate height for the next bullet
        chars = sum(len(t) for t, _, _ in runs)
        per_line = int(w.inches * 100 / (size * 0.062) / 10)
        n = max(1, -(-chars // max(20, per_line)))
        cur += Inches(0.02) + Emu(int(n * size * 1.42 * 12700)) + Emu(int(gap * 12700))
    return cur


def table(slide, x, y, w, rows, widths, size=13.5, head=True, rh=0.38):
    """rows[0] = header row."""
    n, m = len(rows), len(rows[0])
    tbl = slide.shapes.add_table(n, m, x, y, w, Inches(rh * n)).table
    total = sum(widths)
    for j, ww in enumerate(widths):
        tbl.columns[j].width = Emu(int(w * ww / total))
    for i, row in enumerate(rows):
        tbl.rows[i].height = Inches(rh if i else rh)
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.text = ""
            cell.margin_left = cell.margin_right = Inches(0.09)
            cell.margin_top = cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            if i == 0 and head:
                cell.fill.fore_color.rgb = NAVY
            else:
                cell.fill.fore_color.rgb = WHITE if i % 2 else LIGHT
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if (j and len(val) < 14) else PP_ALIGN.LEFT
            r = p.add_run(); r.text = val
            r.font.size = Pt(size)
            r.font.name = MONO if val.startswith(("push", "pop", "top_", "count_", "O(")) else BODY
            r.font.bold = (i == 0 and head)
            r.font.color.rgb = WHITE if (i == 0 and head) else DARK
    return tbl


def footer(slide, n):
    text(slide, Inches(0.55), Inches(7.02), Inches(8.0), Inches(0.3),
         "Hostel Student Management System  |  Array, Stack, String  |  C++",
         size=11, color=GREY)
    text(slide, Inches(12.2), Inches(7.02), Inches(0.6), Inches(0.3),
         str(n), size=11, color=GREY, align=PP_ALIGN.RIGHT)


# =====================================================================
# 1. TITLE
# =====================================================================
s = blank()
box(s, 0, 0, W, H, fill=NAVY)
box(s, 0, Inches(4.62), W, Inches(0.055), fill=BLUE)
text(s, Inches(1.0), Inches(2.05), Inches(11.5), Inches(1.0),
     "Hostel Student Management System", size=41, color=WHITE, bold=True)
text(s, Inches(1.0), Inches(3.15), Inches(11.3), Inches(0.6),
     [[("A Data Structures mini project in ", False, RGBColor(0xC3, 0xD2, 0xF0)),
       ("C++", True, WHITE)]], size=22)
text(s, Inches(1.0), Inches(3.72), Inches(11.3), Inches(0.5),
     [[("Arrays", True, WHITE), ("   •   ", False, BLUE),
       ("Stacks", True, WHITE), ("   •   ", False, BLUE),
       ("Strings", True, WHITE)]], size=21)
text(s, Inches(1.0), Inches(5.05), Inches(11.3), Inches(1.2),
     ["DS Project 2",
      "Console menu  +  web interface (HTML & CSS) on localhost",
      "No external libraries — every data structure written by hand"],
     size=15, color=RGBColor(0xA9, 0xBE, 0xE4), space=4)

# =====================================================================
# 2. OBJECTIVE
# =====================================================================
s = blank(); header(s, "What the Project Does", "01")
text(s, Inches(0.55), Inches(1.32), Inches(12.2), Inches(0.6),
     [[("Goal: ", True, NAVY),
       ("manage hostel student records while implementing the basic data "
        "structures by hand, instead of using ready-made containers from the "
        "C++ standard library.", False, DARK)]], size=17, line=1.25)

box(s, Inches(0.55), Inches(2.35), Inches(5.9), Inches(3.5), fill=LIGHT, line=BORDER)
text(s, Inches(0.85), Inches(2.58), Inches(5.3), Inches(0.4),
     "FEATURES", size=13, color=BLUE, bold=True)
bullets(s, Inches(0.85), Inches(3.05), Inches(5.3), [
    ("Add", " a student record"),
    ("Display", " all students"),
    ("Search", " by name, course or roll number"),
    ("Delete", " a student"),
    ("Undo", " the last deletion"),
    ("Sort", " alphabetically by name"),
], size=15.5, gap=7)

box(s, Inches(6.85), Inches(2.35), Inches(5.9), Inches(3.5), fill=WHITE, line=BLUE)
text(s, Inches(7.15), Inches(2.58), Inches(5.3), Inches(0.4),
     "ONE RECORD", size=13, color=BLUE, bold=True)
code(s, Inches(7.15), Inches(3.05), Inches(5.3), [
    "struct Student {",
    "    int    roll;",
    "    string name;      // STRING",
    "    string course;    // STRING",
    "    int    room;",
    "};",
], size=15)
text(s, Inches(7.15), Inches(5.05), Inches(5.3), Inches(0.7),
     "100 such records are kept in one array, and deleted ones are kept on a stack.",
     size=14, color=GREY, line=1.2)
footer(s, 2)

# =====================================================================
# 3. WHY THESE THREE STRUCTURES
# =====================================================================
s = blank(); header(s, "Why These Three Structures", "02")
text(s, Inches(0.55), Inches(1.3), Inches(12.2), Inches(0.4),
     "Each structure was chosen for a reason — not just to tick a requirement.",
     size=16, color=GREY)
table(s, Inches(0.55), Inches(1.85), Inches(12.2), [
    ["Structure", "Where it is used", "Why this one"],
    ["ARRAY",
     "students[100] holds every record, with count_ tracking how many are used",
     "Records need indexed access and ordered display, and the maximum is known in advance"],
    ["STACK",
     "Deleted records are pushed here; Undo pops the top one back",
     "Undo must reverse the most recent action first — that is exactly LIFO"],
    ["STRING",
     "Name and course fields; case-insensitive search and alphabetical sort",
     "Text data needs character-level comparison, which std::string supports directly"],
], [1.6, 5.0, 5.6], size=14.5, rh=0.92)

box(s, Inches(0.55), Inches(5.72), Inches(12.2), Inches(1.05), fill=RGBColor(0xEC, 0xFD, 0xF3), line=GREEN)
text(s, Inches(0.85), Inches(5.95), Inches(11.6), Inches(0.7),
     [[("The key point: ", True, RGBColor(0x0F, 0x7A, 0x38)),
       ("the stack is not decoration. Undo has to reverse the ", False, DARK),
       ("most recent", True, DARK),
       (" deletion first, which is Last In First Out. A queue would be wrong — "
        "it would undo the ", False, DARK), ("oldest", True, DARK),
       (" deletion first.", False, DARK)]], size=16, line=1.22)
footer(s, 3)

# =====================================================================
# 4. THE ARRAY
# =====================================================================
s = blank(); header(s, "The Array — Storing the Records", "03")
code(s, Inches(0.55), Inches(1.3), Inches(5.85), [
    "Student students[MAX];   // MAX = 100",
    "int count_ = 0;          // how many are used",
    "",
    "bool addStudent(int roll, string name,",
    "                string course, int room) {",
    "  if (count_ >= MAX)   return false;  // full",
    "  if (findStudent(roll) != -1)",
    "      return false;    // duplicate roll",
    "",
    "  students[count_].roll   = roll;",
    "  students[count_].name   = name;",
    "  students[count_].course = course;",
    "  students[count_].room   = room;",
    "  count_++;",
    "  return true;",
    "}",
], size=13.5)

text(s, Inches(6.85), Inches(1.3), Inches(5.9), Inches(0.4),
     "HOW IT WORKS", size=13, color=BLUE, bold=True)
bullets(s, Inches(6.85), Inches(1.78), Inches(5.9), [
    ("Insert ", "puts the record at index count_ and increments it — O(1), no shifting needed"),
    ("count_ ", "is the logical size, kept separate from the array's physical capacity of 100"),
    ("Two checks first: ", "the array must not be full, and the roll number must not already exist"),
    ("Search ", "walks index 0 to count_-1 — a linear search, O(n)"),
    ("findStudent() returns the index", ", not the record, because deletion needs to know where it is"),
], size=15.5, gap=9)

box(s, Inches(6.85), Inches(5.5), Inches(5.9), Inches(0.95), fill=LIGHT, line=BORDER)
code(s, Inches(7.05), Inches(5.65), Inches(5.5), [
    "for (int i = 0; i < count_; i++)      // linear search",
    "    if (students[i].roll == roll) return i;",
], size=12.5)
footer(s, 4)

# =====================================================================
# 5. THE STACK
# =====================================================================
s = blank(); header(s, "The Stack — Remembering Deletions", "04")
code(s, Inches(0.55), Inches(1.3), Inches(5.85), [
    "Student stack_[MAX];   // the stack",
    "int top_ = -1;         // -1 means empty",
    "",
    "void push(Student s) {",
    "    if (top_ < MAX - 1) {   // overflow check",
    "        top_++;",
    "        stack_[top_] = s;",
    "    }",
    "}",
    "",
    "Student pop() {",
    "    Student s = stack_[top_];",
    "    top_--;",
    "    return s;",
    "}",
], size=13.5, accent=GREEN)

text(s, Inches(6.85), Inches(1.3), Inches(5.9), Inches(0.4),
     "ARRAY IMPLEMENTATION OF A STACK", size=13, color=BLUE, bold=True)
table(s, Inches(6.85), Inches(1.75), Inches(5.9), [
    ["Operation", "What it does", "Time"],
    ["push(s)", "top_++ then store at that index", "O(1)"],
    ["pop()", "read at top_ then top_--", "O(1)"],
    ["stackEmpty()", "true when top_ == -1", "O(1)"],
], [1.5, 3.3, 0.9], size=13.5, rh=0.46)

bullets(s, Inches(6.85), Inches(3.75), Inches(5.9), [
    ("Overflow ", "= pushing onto a full stack, blocked by top_ < MAX-1"),
    ("Underflow ", "= popping an empty stack, blocked by top_ == -1"),
    ("Both checked before ", "the array is touched, so memory outside it is never accessed"),
    ("top_ starts at -1 ", "because an empty stack has no valid top index; size is always top_ + 1"),
], size=15.5, gap=9)
footer(s, 5)

# =====================================================================
# 6. STRINGS
# =====================================================================
s = blank(); header(s, "Strings — Searching and Sorting", "05")

text(s, Inches(0.55), Inches(1.3), Inches(5.85), Inches(0.4),
     "CASE-INSENSITIVE SEARCH", size=13, color=BLUE, bold=True)
code(s, Inches(0.55), Inches(1.75), Inches(5.85), [
    "string toLower(string s) {",
    "  for (int i = 0; i < s.size(); i++)",
    "    if (s[i] >= 'A' && s[i] <= 'Z')",
    "        s[i] = s[i] + 32;   // 'A'=65 'a'=97",
    "  return s;",
    "}",
    "",
    "bool contains(string text, string key) {",
    "  return toLower(text).find(toLower(key))",
    "         != string::npos;",
    "}",
], size=13.5)
box(s, Inches(0.55), Inches(4.95), Inches(5.85), Inches(0.75), fill=LIGHT, line=BORDER)
text(s, Inches(0.78), Inches(5.12), Inches(5.5), Inches(0.5),
     [[("This is why typing ", False, DARK), ("bca", True, NAVY),
       (" also finds ", False, DARK), ("BCA", True, NAVY),
       (". The stored data is never changed — only the copies used for comparison.", False, DARK)]],
     size=14, line=1.2)

text(s, Inches(6.9), Inches(1.3), Inches(5.85), Inches(0.4),
     "SORTING BY NAME (BUBBLE SORT)", size=13, color=BLUE, bold=True)
code(s, Inches(6.9), Inches(1.75), Inches(5.85), [
    "void sortByName() {",
    " for (int i = 0; i < count_-1; i++)",
    "  for (int j = 0; j < count_-1-i; j++)",
    "   if (toLower(students[j].name) >",
    "       toLower(students[j+1].name)) {",
    "      Student t     = students[j];",
    "      students[j]   = students[j+1];",
    "      students[j+1] = t;      // swap",
    "   }",
    "}",
], size=13.5)
box(s, Inches(6.9), Inches(4.72), Inches(5.85), Inches(0.98), fill=LIGHT, line=BORDER)
text(s, Inches(7.13), Inches(4.9), Inches(5.5), Inches(0.7),
     [[("The ", False, DARK), (">", True, NAVY),
       (" operator on std::string compares ", False, DARK),
       ("lexicographically", True, NAVY),
       (" — character by character, like a dictionary. Nested loops make this ", False, DARK),
       ("O(n²)", True, NAVY), (".", False, DARK)]], size=14, line=1.2)

text(s, Inches(0.55), Inches(6.0), Inches(12.2), Inches(0.5),
     [[("Search also matches the roll number typed as text, using ", False, DARK),
       ("to_string()", True, NAVY),
       (" — so one search box covers name, course and roll.", False, DARK)]],
     size=15, color=GREY)
footer(s, 6)

# =====================================================================
# 7. DELETE = the key slide
# =====================================================================
s = blank(); header(s, "Deletion: Both Structures Working Together", "06")
text(s, Inches(0.55), Inches(1.25), Inches(12.2), Inches(0.4),
     [[("This is the most important function in the project — the array and the "
        "stack are used in the ", False, DARK), ("same operation", True, NAVY),
       (".", False, DARK)]], size=16.5)

code(s, Inches(0.55), Inches(1.85), Inches(6.4), [
    "bool deleteStudent(int roll) {",
    "  int i = findStudent(roll);",
    "  if (i == -1) return false;   // not found",
    "",
    "  push(students[i]);      // STACK: save first",
    "",
    "  for (int j = i; j < count_-1; j++)",
    "      students[j] = students[j+1];  // shift",
    "  count_--;",
    "  return true;",
    "}",
], size=14, accent=RED)

box(s, Inches(0.55), Inches(4.85), Inches(6.4), Inches(1.5), fill=RGBColor(0xFF, 0xF7, 0xE6), line=RGBColor(0xE0, 0xA8, 0x00))
text(s, Inches(0.8), Inches(5.05), Inches(5.95), Inches(1.15),
     [[("The order matters. ", True, RGBColor(0x8A, 0x5A, 0x00)),
       ("The record is pushed onto the stack ", False, DARK),
       ("before", True, DARK),
       (" the shift overwrites it. Push after shifting would copy the wrong "
        "student, and undo would restore garbage.", False, DARK)]], size=15.5, line=1.25)

text(s, Inches(7.35), Inches(1.85), Inches(5.4), Inches(0.4),
     "DELETING ROLL 102", size=13, color=BLUE, bold=True)
box(s, Inches(7.35), Inches(2.3), Inches(5.4), Inches(4.05), fill=CODEBG)
tb = s.shapes.add_textbox(Inches(7.55), Inches(2.48), Inches(5.05), Inches(3.7))
tf = tb.text_frame; tf.word_wrap = False
tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
diagram = [
    ("before:  [101][102][103][104]", CODEFG),
    ("                |             ", CODEFG),
    ("         count_ = 4           ", RGBColor(0x8A, 0x93, 0xA5)),
    ("", CODEFG),
    ("step 1   push a copy onto the stack", RGBColor(0x7E, 0xC9, 0x9B)),
    ("", CODEFG),
    ("shift:   [101][103][104][104]", CODEFG),
    ("          move each one left  ", RGBColor(0x8A, 0x93, 0xA5)),
    ("", CODEFG),
    ("after:   [101][103][104]", CODEFG),
    ("         count_ = 3           ", RGBColor(0x8A, 0x93, 0xA5)),
    ("         last slot is ignored ", RGBColor(0x8A, 0x93, 0xA5)),
    ("", CODEFG),
    ("stack:   | 102 |  <-- top_", RGBColor(0xF0, 0xB4, 0x5C)),
]
for i, (ln, col) in enumerate(diagram):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(0); p.line_spacing = 1.18
    r = p.add_run(); r.text = ln
    r.font.size = Pt(14); r.font.name = MONO; r.font.color.rgb = col
footer(s, 7)

# =====================================================================
# 8. UNDO
# =====================================================================
s = blank(); header(s, "Undo — Popping the Stack", "07")
code(s, Inches(0.55), Inches(1.35), Inches(6.4), [
    "bool undoDelete() {",
    "  if (stackEmpty()) return false;  // nothing",
    "  Student s = pop();     // STACK: top record",
    "  return addStudent(s.roll, s.name,",
    "                    s.course, s.room);",
    "}",
], size=14, accent=GREEN)

bullets(s, Inches(0.55), Inches(3.35), Inches(6.4), [
    ("Reuses addStudent()", " — the restored record goes back through the normal insertion path"),
    ("Underflow is checked first", ", so pressing Undo with an empty stack is safe"),
    ("LIFO in action: ", "delete three students, press Undo three times, and they come back in reverse order"),
], size=16, gap=10)

text(s, Inches(7.35), Inches(1.35), Inches(5.4), Inches(0.4),
     "WHY LIFO IS THE RIGHT BEHAVIOUR", size=13, color=BLUE, bold=True)
box(s, Inches(7.35), Inches(1.8), Inches(5.4), Inches(3.6), fill=LIGHT, line=BORDER)
tb = s.shapes.add_textbox(Inches(7.6), Inches(2.0), Inches(4.95), Inches(3.25))
tf = tb.text_frame; tf.word_wrap = False
tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
seq = [
    ("delete 101   stack:  | 101 |", NAVY),
    ("delete 102   stack:  | 101 | 102 |  <- top", NAVY),
    ("delete 103   stack:  | 101 | 102 | 103 |  <- top", NAVY),
    ("", DARK),
    ("undo   ->  103 comes back  (newest first)", RGBColor(0x0F, 0x7A, 0x38)),
    ("undo   ->  102 comes back", RGBColor(0x0F, 0x7A, 0x38)),
    ("undo   ->  101 comes back", RGBColor(0x0F, 0x7A, 0x38)),
    ("undo   ->  stack empty, nothing happens", GREY),
]
for i, (ln, col) in enumerate(seq):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(4); p.line_spacing = 1.15
    r = p.add_run(); r.text = ln
    r.font.size = Pt(13.5); r.font.name = MONO; r.font.color.rgb = col

text(s, Inches(7.35), Inches(5.6), Inches(5.4), Inches(0.9),
     "A queue here would undo 101 first — the oldest deletion — which is not what "
     "undo means. That is why the structure had to be a stack.",
     size=15, color=GREY, line=1.22)
footer(s, 8)

# =====================================================================
# 9. STRUCTURE
# =====================================================================
s = blank(); header(s, "Project Structure", "08")
box(s, Inches(0.55), Inches(1.3), Inches(6.6), Inches(4.5), fill=CODEBG)
tb = s.shapes.add_textbox(Inches(0.8), Inches(1.52), Inches(6.2), Inches(4.1))
tf = tb.text_frame; tf.word_wrap = False
tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
tree = [
    ("ds_project_2/", CODEFG),
    (" |", GREY),
    " +-- src/",
    ("      hostel.h      ALL data structure code", RGBColor(0xF0, 0xB4, 0x5C)),
    "      console.cpp   the console menu program",
    "      server.cpp     the small web server",
    (" |", GREY),
    " +-- web/",
    "      index.html    the interface (no JavaScript)",
    "      style.css     the styling",
    (" |", GREY),
    " +-- docs/          report and this presentation",
    " +-- Makefile       build commands",
    " +-- README.md",
]
for i, ln in enumerate(tree):
    col = CODEFG
    if isinstance(ln, tuple):
        ln, col = ln
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(1); p.line_spacing = 1.14
    r = p.add_run(); r.text = ln
    r.font.size = Pt(14); r.font.name = MONO; r.font.color.rgb = col

text(s, Inches(7.55), Inches(1.3), Inches(5.2), Inches(0.4),
     "ONE FILE HOLDS THE LOGIC", size=13, color=BLUE, bold=True)
text(s, Inches(7.55), Inches(1.78), Inches(5.2), Inches(1.6),
     [[("All the data structures live in ", False, DARK),
       ("src/hostel.h", True, NAVY),
       (" — about 140 lines. The console program and the web server contain ", False, DARK),
       ("no data structure code of their own", True, NAVY),
       (": they are just two different ways of calling the same functions.", False, DARK)]],
     size=16, line=1.28)
bullets(s, Inches(7.55), Inches(3.65), Inches(5.2), [
    ("Separation of logic from interface", ""),
    ("The same header can be dropped into a bigger project", ""),
    ("Nothing to install — standard library only", ""),
], size=15.5, gap=9)
footer(s, 9)

# =====================================================================
# 10. HOW TO RUN
# =====================================================================
s = blank(); header(s, "How to Run It", "09")
text(s, Inches(0.55), Inches(1.28), Inches(6.0), Inches(0.4),
     "WEB VERSION  (use this for the demo)", size=13, color=BLUE, bold=True)
code(s, Inches(0.55), Inches(1.72), Inches(6.0), [
    "$ cd ds_project_2",
    "$ make web",
    "",
    "g++ src/server.cpp -o server",
    "./server",
    "Server running -> http://localhost:8080",
], size=14, accent=GREEN)
text(s, Inches(0.55), Inches(3.6), Inches(6.0), Inches(0.4),
     "CONSOLE VERSION", size=13, color=BLUE, bold=True)
code(s, Inches(0.55), Inches(4.04), Inches(6.0), [
    "$ make cli",
    "",
    "1. Add student      5. Undo delete (stack)",
    "2. Show all         6. Sort by name",
    "3. Search student   7. Exit",
    "4. Delete student",
], size=14, accent=GREEN)

text(s, Inches(6.95), Inches(1.28), Inches(5.8), Inches(0.4),
     "REQUIREMENTS", size=13, color=BLUE, bold=True)
bullets(s, Inches(6.95), Inches(1.72), Inches(5.8), [
    ("g++ ", "and Linux or macOS — nothing to install"),
    ("POSIX sockets ", "are built into the operating system"),
    ("Any browser ", "for the web version"),
], size=15.5, gap=8)

box(s, Inches(6.95), Inches(3.5), Inches(5.8), Inches(1.35), fill=RGBColor(0xFF, 0xF7, 0xE6), line=RGBColor(0xE0, 0xA8, 0x00))
text(s, Inches(7.2), Inches(3.72), Inches(5.3), Inches(1.0),
     [[("Run it from inside the project folder. ", True, RGBColor(0x8A, 0x5A, 0x00)),
       ("The server reads web/index.html and web/style.css by relative path, so "
        "starting it from elsewhere gives a blank page.", False, DARK)]],
     size=15, line=1.22)
text(s, Inches(6.95), Inches(5.1), Inches(5.8), Inches(1.0),
     [[("Other commands: ", False, GREY), ("make", True, NAVY),
       (" builds both programs, ", False, GREY), ("make clean", True, NAVY),
       (" removes them, ", False, GREY), ("make report", True, NAVY),
       (" rebuilds the PDF.", False, GREY)]], size=15, line=1.22)
footer(s, 10)

# =====================================================================
# 11. WEB UI
# =====================================================================
s = blank(); header(s, "How the Web Interface Works", "10")
text(s, Inches(0.55), Inches(1.25), Inches(12.2), Inches(0.4),
     [[("The interface is ", False, DARK), ("HTML and CSS only", True, NAVY),
       (" — there is no JavaScript. Every button is a normal form or link, and "
        "C++ does all the work.", False, DARK)]], size=16.5)

box(s, Inches(0.55), Inches(1.85), Inches(7.6), Inches(4.35), fill=CODEBG)
tb = s.shapes.add_textbox(Inches(0.78), Inches(2.05), Inches(7.2), Inches(4.0))
tf = tb.text_frame; tf.word_wrap = False
tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
flow = [
    ("  browser                     C++ server (server.cpp)", RGBColor(0x9F, 0xB6, 0xE8)),
    ("     |                              |", GREY),
    ("  GET /  ---------------------->     read web/index.html", CODEFG),
    ("     |                              loop the ARRAY -> table rows", RGBColor(0x7E, 0xC9, 0x9B)),
    ("     |                              loop the STACK -> the list", RGBColor(0x7E, 0xC9, 0x9B)),
    ("     |                              fill in {{ROWS}} {{STACK}}", RGBColor(0x7E, 0xC9, 0x9B)),
    ("     |  <--- finished HTML page      {{COUNT}} {{SEARCH}}", CODEFG),
    ("     |                              |", GREY),
    ("  POST /add  ------------------>     addStudent()   array insert", CODEFG),
    ("  GET /delete?roll=102 ------->     deleteStudent() shift + push", CODEFG),
    ("  GET /undo  ------------------>     undoDelete()   pop", CODEFG),
    ("  GET /sort  ------------------>     sortByName()   string compare", CODEFG),
    ("  GET /?q=bca ----------------->     matches()      string search", CODEFG),
    ("     |  <--- redirect back to /      |", GREY),
]
for i, (ln, col) in enumerate(flow):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(1); p.line_spacing = 1.16
    r = p.add_run(); r.text = ln
    r.font.size = Pt(13); r.font.name = MONO; r.font.color.rgb = col

text(s, Inches(8.55), Inches(1.85), Inches(4.2), Inches(0.4),
     "TEMPLATE REPLACEMENT", size=13, color=BLUE, bold=True)
text(s, Inches(8.55), Inches(2.3), Inches(4.2), Inches(1.5),
     [[("index.html holds markers like ", False, DARK), ("{{ROWS}}", True, NAVY),
       (". The C++ function makePage() swaps each one for text built from the "
        "array and the stack — so design stays separate from logic.", False, DARK)]],
     size=15, line=1.25)
box(s, Inches(8.55), Inches(4.15), Inches(4.2), Inches(2.05), fill=RGBColor(0xEC, 0xFD, 0xF3), line=GREEN)
text(s, Inches(8.78), Inches(4.35), Inches(3.75), Inches(1.7),
     [[("Good for the demo: ", True, RGBColor(0x0F, 0x7A, 0x38)),
       ("the page shows the stack itself, with the top marked ", False, DARK),
       ("TOP →", True, NAVY),
       (". Delete a student and the record appears on the stack; press Undo and "
        "it disappears. The structure is visible on screen.", False, DARK)]],
     size=14.5, line=1.22)
footer(s, 11)

# =====================================================================
# 12. TESTING
# =====================================================================
s = blank(); header(s, "Testing", "11")
text(s, Inches(0.55), Inches(1.25), Inches(12.2), Inches(0.4),
     "Every feature was tested in both versions — these are the actual results.",
     size=16, color=GREY)
table(s, Inches(0.55), Inches(1.8), Inches(12.2), [
    ["#", "Test", "Expected result", "Status"],
    ["1", "Add a student with a new roll number", "Added, array size grows by 1", "Pass"],
    ["2", "Add a roll number that already exists", "Rejected, array unchanged", "Pass"],
    ["3", "Search  bca  in small letters", "Finds the student whose course is BCA", "Pass"],
    ["4", "Delete a student from the middle", "Removed, later records shift left, no gap", "Pass"],
    ["5", "Check the stack after deleting", "Deleted record shown on top of the stack", "Pass"],
    ["6", "Press Undo", "Student returns, stack becomes empty", "Pass"],
    ["7", "Press Undo with an empty stack", "Handled safely, no crash (underflow check)", "Pass"],
    ["8", "Delete two students, then undo twice", "Restored in reverse order (LIFO)", "Pass"],
    ["9", "Sort by name", "Anita, Karan, Priya, Ravi — alphabetical", "Pass"],
    ["10", "Add a name with a space via the web form", "Stored correctly as \"Priya Singh\"", "Pass"],
], [0.6, 4.6, 5.6, 1.1], size=13.5, rh=0.41)
text(s, Inches(0.55), Inches(6.5), Inches(12.2), Inches(0.4),
     "12 test cases in total (the full table is in the project report) — all passing.",
     size=14.5, color=GREY)
footer(s, 12)

# =====================================================================
# 13. COMPLEXITY
# =====================================================================
s = blank(); header(s, "Time Complexity", "12")
table(s, Inches(0.55), Inches(1.4), Inches(12.2), [
    ["Operation", "Structures used", "Best", "Worst", "Reason"],
    ["Add student", "Array", "O(1)", "O(n)", "insert is O(1), but the duplicate check is a search"],
    ["Search", "Array + String", "O(1)", "O(n)", "linear search — the whole array in the worst case"],
    ["Delete student", "Array + Stack", "O(1)", "O(n)", "shifting every element after the deleted one"],
    ["Push / Pop", "Stack", "O(1)", "O(1)", "only the top_ index changes"],
    ["Undo delete", "Stack + Array", "O(1)", "O(n)", "O(1) pop, then a re-insert that checks duplicates"],
    ["Sort by name", "Array + String", "O(n)", "O(n²)", "bubble sort with nested loops"],
], [2.3, 2.4, 1.0, 1.0, 5.5], size=13.5, rh=0.55)

box(s, Inches(0.55), Inches(5.62), Inches(5.9), Inches(1.1), fill=LIGHT, line=BORDER)
text(s, Inches(0.8), Inches(5.84), Inches(5.45), Inches(0.8),
     [[("Space complexity: ", True, NAVY),
       ("O(n) — one array of 100 records plus one stack of 100 records, "
        "allocated once at the start.", False, DARK)]], size=15.5, line=1.22)
box(s, Inches(6.85), Inches(5.62), Inches(5.9), Inches(1.1), fill=LIGHT, line=BORDER)
text(s, Inches(7.1), Inches(5.84), Inches(5.45), Inches(0.8),
     [[("Note: ", True, NAVY),
       ("push and pop are the only O(1) worst case operations — that is the "
        "advantage of using a stack for undo.", False, DARK)]], size=15.5, line=1.22)
footer(s, 13)

# =====================================================================
# 14. LIMITATIONS
# =====================================================================
s = blank(); header(s, "Limitations and Future Scope", "13")
table(s, Inches(0.55), Inches(1.35), Inches(12.2), [
    ["Current limitation", "How it could be improved"],
    ["Records are lost when the program stops (memory only)", "Save to a text or CSV file using file handling"],
    ["Fixed capacity of 100 students", "A dynamic array that doubles in size, or a linked list"],
    ["Search is a linear scan, O(n)", "Keep the array sorted by roll number and use binary search, O(log n)"],
    ["Bubble sort is O(n²)", "Merge sort or quick sort, O(n log n)"],
    ["Undo covers deletion only", "Push every action onto the stack to undo edits as well"],
    ["The server handles one request at a time", "Handle each connection in a separate thread"],
    ["No room allocation or waiting list", "Add a queue for the waiting list — naturally FIFO"],
], [5.9, 6.3], size=14.5, rh=0.57)
text(s, Inches(0.55), Inches(6.55), Inches(12.2), Inches(0.5),
     [[("The natural next step: ", True, NAVY),
       ("a queue for the room waiting list. Students should get rooms in the "
        "order they applied — First In First Out — just as undo needed Last In First Out.", False, DARK)]],
     size=15.5, line=1.2)
footer(s, 14)

# =====================================================================
# 15. VIVA
# =====================================================================
s = blank(); header(s, "Questions You May Be Asked", "14")
qa = [
    ("Why a stack and not a queue?",
     "Undo must reverse the most recent deletion first — LIFO. A queue would undo the oldest first."),
    ("Why not use std::vector or std::stack?",
     "The point is to demonstrate the structures, so both are written by hand: count_ for the array, top_ for the stack."),
    ("What happens when you delete from the middle?",
     "Linear search finds the position, the record is pushed, then everything after it shifts one place left. O(n)."),
    ("What are overflow and underflow?",
     "Pushing onto a full stack and popping an empty one. Both are checked before the array is touched."),
    ("Why does top_ start at -1?",
     "An empty stack has no valid top index, and -1 is just before index 0. It also makes size = top_ + 1."),
    ("How does bca match BCA?",
     "Both strings are converted with toLower(), which adds 32 to the ASCII code of capitals, before find() compares them."),
]
y = Inches(1.28)
for i, (q, a) in enumerate(qa):
    box(s, Inches(0.55), y, Inches(12.2), Inches(0.86),
        fill=LIGHT if i % 2 == 0 else WHITE, line=BORDER)
    text(s, Inches(0.8), y + Inches(0.1), Inches(11.7), Inches(0.3),
         "Q.  " + q, size=15.5, color=NAVY, bold=True)
    text(s, Inches(0.8), y + Inches(0.42), Inches(11.7), Inches(0.36),
         "A.  " + a, size=14.5, color=DARK)
    y += Inches(0.9)
text(s, Inches(0.55), Inches(6.85), Inches(12.2), Inches(0.4),
     "Ten questions with full answers are in section 9 of the project report.",
     size=14, color=GREY)
footer(s, 15)

# =====================================================================
# 16. CONCLUSION
# =====================================================================
s = blank()
box(s, 0, 0, W, H, fill=NAVY)
box(s, Inches(1.0), Inches(1.15), Inches(0.055), Inches(0.85), fill=BLUE)
text(s, Inches(1.3), Inches(1.2), Inches(11.0), Inches(0.8),
     "Conclusion", size=38, color=WHITE, bold=True)
pts = [
    ("Array", " — indexed storage of the records, with insert, linear search and delete-with-shift"),
    ("Stack", " — the undo history, chosen because undo is inherently Last In First Out"),
    ("String", " — case-insensitive search and alphabetical sorting through character comparison"),
]
y = Inches(2.5)
for lead, rest in pts:
    box(s, Inches(1.3), y + Inches(0.12), Inches(0.08), Inches(0.08), fill=BLUE)
    text(s, Inches(1.62), y, Inches(10.5), Inches(0.5),
         [[(lead, True, WHITE), (rest, False, RGBColor(0xC3, 0xD2, 0xF0))]],
         size=19, line=1.2)
    y += Inches(0.72)

text(s, Inches(1.3), Inches(4.9), Inches(10.8), Inches(1.2),
     [[("All three are implemented from first principles, and the same logic is "
        "exposed through two interfaces — a console menu and a browser page — "
        "which shows that the data structure code is independent of how it is "
        "presented.", False, RGBColor(0xC3, 0xD2, 0xF0))]], size=17, line=1.3)
box(s, Inches(1.3), Inches(6.15), Inches(10.8), Inches(0.02), fill=BLUE)
text(s, Inches(1.3), Inches(6.35), Inches(10.8), Inches(0.5),
     [[("Everything lives in one header, ", False, RGBColor(0x9F, 0xB6, 0xE8)),
       ("src/hostel.h", True, WHITE),
       (" — ready to reuse in a larger project.", False, RGBColor(0x9F, 0xB6, 0xE8))]],
     size=16)

prs.save("docs/DS_Project_2_Presentation.pptx")
print("saved docs/DS_Project_2_Presentation.pptx with %d slides" % len(prs.slides._sldIdLst))
